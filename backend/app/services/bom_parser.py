import uuid
import io
from typing import IO
from app.models.bom import BomItem, BomImportResult

# 장비별 기본 스펙 (알려진 장비는 자동 매핑)
KNOWN_SPECS: dict[str, dict] = {
    "hgx b300 server":   {"unit_power_w": 10000, "unit_height_u": 10},
    "rtx server":        {"unit_power_w": 3000,  "unit_height_u": 4},
    "mgmt server":       {"unit_power_w": 400,   "unit_height_u": 2},
    "nas storage":       {"unit_power_w": 500,   "unit_height_u": 2},
    "firewall":          {"unit_power_w": 200,   "unit_height_u": 1},
    "ai switch":         {"unit_power_w": 1500,  "unit_height_u": 2},
    "leaf switch":       {"unit_power_w": 500,   "unit_height_u": 1},
    "mgmt. switch":      {"unit_power_w": 200,   "unit_height_u": 1},
    "in-row cooling":    {"unit_power_w": None,  "unit_height_u": 42},
    "cdu":               {"unit_power_w": None,  "unit_height_u": 4},
}


def _lookup_specs(equipment: str) -> dict:
    key = equipment.lower().strip()
    for k, v in KNOWN_SPECS.items():
        if k in key:
            return v
    return {}


def _fill_merged_cells(rows: list[list[str]]) -> list[list[str]]:
    """PPTX 병합 셀(빈 값)을 위 행 값으로 채움."""
    result = []
    prev = [""] * (len(rows[0]) if rows else 5)
    for row in rows:
        filled = []
        for i, cell in enumerate(row):
            filled.append(cell if cell else prev[i])
        prev = filled
        result.append(filled)
    return result


def parse_pptx(file: IO[bytes]) -> BomImportResult:
    from pptx import Presentation
    prs = Presentation(file)
    warnings: list[str] = []
    raw_rows: list[list[str]] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                for r in range(len(tbl.rows)):
                    raw_rows.append(
                        [tbl.cell(r, c).text.strip() for c in range(len(tbl.columns))]
                    )
                break  # 첫 번째 테이블만

    if not raw_rows:
        return BomImportResult(items=[], source="pptx", warnings=["테이블을 찾을 수 없습니다."])

    # 헤더 행 제거
    header = raw_rows[0]
    data_rows = _fill_merged_cells(raw_rows[1:])

    items: list[BomItem] = []
    for row in data_rows:
        if len(row) < 4:
            continue
        category, subcategory, equipment, qty_str, *rest = row
        description = rest[0] if rest else ""
        try:
            qty = int(qty_str)
        except ValueError:
            warnings.append(f"수량 파싱 실패: {row}")
            qty = 0

        specs = _lookup_specs(equipment)
        items.append(BomItem(
            id=str(uuid.uuid4()),
            category=category,
            subcategory=subcategory,
            equipment=equipment,
            quantity=qty,
            description=description,
            **specs,
        ))

    return BomImportResult(items=items, source="pptx", warnings=warnings)


def parse_excel(file: IO[bytes]) -> BomImportResult:
    import openpyxl
    warnings: list[str] = []
    wb = openpyxl.load_workbook(file, data_only=True)
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))

    if not rows:
        return BomImportResult(items=[], source="excel", warnings=["시트가 비어있습니다."])

    # 헤더 자동 감지: 첫 행
    header = [str(c).lower().strip() if c else "" for c in rows[0]]
    data = rows[1:]

    def col(h_keywords: list[str]) -> int:
        for kw in h_keywords:
            for i, h in enumerate(header):
                if kw in h:
                    return i
        return -1

    cat_i   = col(["분야", "category"])
    sub_i   = col(["항목", "subcategory", "sub"])
    eq_i    = col(["equipment", "장비", "모델"])
    qty_i   = col(["qty", "수량", "quantity"])
    desc_i  = col(["설명", "description", "desc"])

    items: list[BomItem] = []
    prev_cat, prev_sub = "", ""
    for row in data:
        def v(i):
            return str(row[i]).strip() if i >= 0 and i < len(row) and row[i] is not None else ""
        category    = v(cat_i) or prev_cat
        subcategory = v(sub_i) or prev_sub
        equipment   = v(eq_i)
        if not equipment:
            continue
        prev_cat, prev_sub = category, subcategory
        try:
            qty = int(float(v(qty_i))) if v(qty_i) else 0
        except ValueError:
            qty = 0
            warnings.append(f"수량 파싱 실패: {row}")
        specs = _lookup_specs(equipment)
        items.append(BomItem(
            id=str(uuid.uuid4()),
            category=category,
            subcategory=subcategory,
            equipment=equipment,
            quantity=qty,
            description=v(desc_i),
            **specs,
        ))

    return BomImportResult(items=items, source="excel", warnings=warnings)


def parse_csv(file: IO[bytes]) -> BomImportResult:
    import csv
    warnings: list[str] = []
    text = file.read().decode("utf-8-sig")
    reader = csv.DictReader(io.StringIO(text))
    items: list[BomItem] = []

    for row in reader:
        def g(keys):
            for k in keys:
                for rk in row:
                    if k.lower() in rk.lower():
                        return row[rk].strip()
            return ""

        equipment = g(["equipment", "장비", "모델"])
        if not equipment:
            continue
        try:
            qty = int(float(g(["qty", "수량"])))
        except ValueError:
            qty = 0
        specs = _lookup_specs(equipment)
        items.append(BomItem(
            id=str(uuid.uuid4()),
            category=g(["분야", "category"]),
            subcategory=g(["항목", "sub"]),
            equipment=equipment,
            quantity=qty,
            description=g(["설명", "description"]),
            **specs,
        ))

    return BomImportResult(items=items, source="csv", warnings=warnings)
